"""
Proposal export services for PDF, DOCX, and PPTX formats.
"""
from typing import List, Dict, Any, Optional
from pathlib import Path
from io import BytesIO
from datetime import datetime
import uuid

# PDF Export
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

# DOCX Export
try:
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# PPTX Export
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

class ProposalExporter:
    """Export proposals to various formats."""
    
    def __init__(self):
        self.export_dir = Path("./exports")
        self.export_dir.mkdir(exist_ok=True)
    
    def export_pdf(
        self,
        title: str,
        sections: List[Dict[str, Any]],
        project_name: str = None,
        client_name: str = None
    ) -> BytesIO:
        """
        Export proposal as PDF.
        
        Args:
            title: Proposal title
            sections: List of section dictionaries
            project_name: Optional project name
            client_name: Optional client name
        
        Returns:
            BytesIO buffer with PDF content
        """
        if not REPORTLAB_AVAILABLE:
            raise ValueError("ReportLab not available. Install with: pip install reportlab")
        
        buffer = BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18
        )
        
        # Styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=HexColor('#0066CC'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=HexColor('#0066CC'),
            spaceAfter=12,
            spaceBefore=12
        )
        
        # Build story
        story = []
        
        # Title
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Header info
        if client_name or project_name:
            header_text = ""
            if client_name:
                header_text += f"<b>Client:</b> {client_name}<br/>"
            if project_name:
                header_text += f"<b>Project:</b> {project_name}<br/>"
            header_text += f"<b>Date:</b> {datetime.now().strftime('%B %d, %Y')}"
            story.append(Paragraph(header_text, styles['Normal']))
            story.append(Spacer(1, 0.3*inch))
        
        # Sections
        for section in sorted(sections, key=lambda x: x.get('order', 0)):
            section_title = section.get('title', '')
            section_content = section.get('content', '')
            
            if section_title:
                story.append(Paragraph(section_title, heading_style))
            
            if section_content:
                # Replace newlines with HTML breaks
                content = section_content.replace('\n', '<br/>')
                story.append(Paragraph(content, styles['Normal']))
                story.append(Spacer(1, 0.2*inch))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        return buffer
    
    def export_docx(
        self,
        title: str,
        sections: List[Dict[str, Any]],
        project_name: str = None,
        client_name: str = None
    ) -> BytesIO:
        """
        Export proposal as DOCX.
        
        Args:
            title: Proposal title
            sections: List of section dictionaries
            project_name: Optional project name
            client_name: Optional client name
        
        Returns:
            BytesIO buffer with DOCX content
        """
        if not DOCX_AVAILABLE:
            raise ValueError("python-docx not available. Install with: pip install python-docx")
        
        doc = Document()
        
        # Title
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Header info
        if client_name or project_name:
            doc.add_paragraph()
            if client_name:
                p = doc.add_paragraph()
                p.add_run('Client: ').bold = True
                p.add_run(client_name)
            if project_name:
                p = doc.add_paragraph()
                p.add_run('Project: ').bold = True
                p.add_run(project_name)
            p = doc.add_paragraph()
            p.add_run('Date: ').bold = True
            p.add_run(datetime.now().strftime('%B %d, %Y'))
            doc.add_paragraph()
        
        # Sections
        for section in sorted(sections, key=lambda x: x.get('order', 0)):
            section_title = section.get('title', '')
            section_content = section.get('content', '')
            
            if section_title:
                doc.add_heading(section_title, level=1)
            
            if section_content:
                # Split by newlines and add paragraphs
                for line in section_content.split('\n'):
                    if line.strip():
                        doc.add_paragraph(line.strip())
                    else:
                        doc.add_paragraph()
        
        # Save to buffer
        buffer = BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        return buffer
    
    def export_pptx(
        self,
        title: str,
        sections: List[Dict[str, Any]],
        project_name: str = None,
        client_name: str = None
    ) -> BytesIO:
        """
        Export proposal as PowerPoint.
        
        Args:
            title: Proposal title
            sections: List of section dictionaries
            project_name: Optional project name
            client_name: Optional client name
        
        Returns:
            BytesIO buffer with PPTX content
        """
        if not PPTX_AVAILABLE:
            raise ValueError("python-pptx not available. Install with: pip install python-pptx")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_text = ""
        if client_name:
            subtitle_text += f"Client: {client_name}\n"
        if project_name:
            subtitle_text += f"Project: {project_name}\n"
        subtitle_text += datetime.now().strftime('%B %d, %Y')
        subtitle_shape.text = subtitle_text
        
        # Section slides
        for section in sorted(sections, key=lambda x: x.get('order', 0)):
            section_title = section.get('title', '')
            section_content = section.get('content', '')
            
            if not section_title and not section_content:
                continue
            
            # Use title and content layout
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = section_title or "Section"
            
            # Add content as bullet points
            tf = body_shape.text_frame
            tf.text = section_content.split('\n')[0] if section_content else ""
            
            # Add remaining lines as bullet points
            for line in section_content.split('\n')[1:]:
                if line.strip():
                    p = tf.add_paragraph()
                    p.text = line.strip()
                    p.level = 0
        
        # Save to buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer
    
    def save_export(
        self,
        buffer: BytesIO,
        format: str,
        proposal_id: int
    ) -> str:
        """
        Save export to file.
        
        Args:
            buffer: File buffer
            format: File format (pdf, docx, pptx)
            proposal_id: Proposal ID
        
        Returns:
            File path
        """
        filename = f"proposal_{proposal_id}_{uuid.uuid4().hex[:8]}.{format}"
        file_path = self.export_dir / filename
        
        with open(file_path, 'wb') as f:
            f.write(buffer.getvalue())
        
        return str(file_path)

# Global instance
proposal_exporter = ProposalExporter()

